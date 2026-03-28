import os
import sys
import argparse
from pathlib import Path
from pptx import Presentation

def obter_arquivos_por_data(pasta: str) -> list[Path]:
    """Retorna os arquivos da pasta ordenados por data de modificação."""
    pasta_path = Path(pasta)
    if not pasta_path.is_dir():
        raise NotADirectoryError(f"A pasta não existe ou não é um diretório: {pasta}")

    arquivos = [f for f in pasta_path.iterdir() if f.is_file()]
    arquivos.sort(key=lambda f: f.stat().st_mtime)
    return arquivos

def carregar_lista_nomes(caminho_arquivo: str) -> list[str]:
    """Carrega a lista de nomes de um arquivo de texto."""
    path = Path(caminho_arquivo)
    if not path.is_file():
        raise FileNotFoundError(f"Arquivo de nomes não encontrado: {caminho_arquivo}")

    with open(path, encoding="utf-8") as f:
        linhas: list[str] = []
        for linha in f:
            texto = linha.strip()
            if not texto or texto.startswith("#") or texto.lower().startswith("python "):
                continue
            linhas.append(texto)
    return linhas

def sanitizar_nome_arquivo(nome: str) -> str:
    """Remove caracteres inválidos para nomes de arquivos no Windows."""
    invalidos = '<>:"/\\|?*'
    resultado = "".join("_" if c in invalidos else c for c in nome)
    return resultado.strip().rstrip(".") or "arquivo"

def extrair_titulos_de_pptx(pasta_caminho: str) -> list[str]:
    """Extrai o título do primeiro slide de cada arquivo .pptx na pasta indicada."""
    arquivos = obter_arquivos_por_data(pasta_caminho)
    lista_titulos = []

    for arquivo in arquivos:
        if arquivo.suffix.lower() == '.pptx':
            try:
                prs = Presentation(arquivo)
                if prs.slides:
                    primeiro_slide = prs.slides[0]
                    if primeiro_slide.shapes.title and primeiro_slide.shapes.title.text:
                        titulo = primeiro_slide.shapes.title.text.strip()
                        titulo_limpo = " ".join(titulo.splitlines())
                        lista_titulos.append(titulo_limpo)
                    else:
                        lista_titulos.append(f"Sem_Titulo_{arquivo.stem}")
            except Exception as e:
                lista_titulos.append(f"Erro_Leitura_{arquivo.stem}")
                
    return lista_titulos

def renomear_arquivos(pasta: str, lista_nomes: list[str], extensao: str | None = None, dry_run: bool = False) -> None:
    """Renomeia os arquivos da pasta usando a lista de nomes fornecida."""
    arquivos = obter_arquivos_por_data(pasta)
    pasta_path = Path(pasta)

    if len(lista_nomes) < len(arquivos):
        print(f"Aviso: {len(arquivos)} arquivos, mas apenas {len(lista_nomes)} nomes.")

    for i, arquivo in enumerate(arquivos):
        if i >= len(lista_nomes):
            break

        nome_novo = sanitizar_nome_arquivo(lista_nomes[i].strip())
        if not nome_novo:
            continue

        if extensao is not None:
            if not extensao.startswith("."):
                extensao = "." + extensao
            nome_final = nome_novo if nome_novo.endswith(extensao) else nome_novo + extensao
        else:
            nome_final = nome_novo if arquivo.suffix and nome_novo.endswith(arquivo.suffix) else nome_novo + arquivo.suffix

        destino = pasta_path / nome_final

        if destino == arquivo or destino.exists():
            continue

        if dry_run:
            print(f"[DRY RUN] {arquivo.name} -> {nome_final}")
        else:
            try:
                arquivo.rename(destino)
                print(f"Renomeado: {arquivo.name} -> {nome_final}")
            except OSError as e:
                print(f"Erro ao renomear {arquivo.name}: {e}")

def main():
    parser = argparse.ArgumentParser(description="Renomeia arquivos de uma pasta por ordem de modificação.")
    parser.add_argument("pasta", help="Caminho da pasta")
    parser.add_argument("nomes", nargs="?", help="Caminho do arquivo .txt com nomes")
    parser.add_argument("--nomes-direto", "-n", nargs="+", help="Nomes na linha de comando")
    parser.add_argument("--extensao", "-e", default=None, help="Extensão a forçar")
    parser.add_argument("--dry-run", action="store_true", help="Modo simulação")
    
    args = parser.parse_args()

    if args.nomes_direto:
        lista_nomes = args.nomes_direto
    elif args.nomes:
        try:
            lista_nomes = carregar_lista_nomes(args.nomes)
        except FileNotFoundError as e:
            print(e, file=sys.stderr)
            sys.exit(1)
    else:
        sys.exit(1)

    try:
        renomear_arquivos(args.pasta, lista_nomes, extensao=args.extensao, dry_run=args.dry_run)
    except NotADirectoryError as e:
        print(e, file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()